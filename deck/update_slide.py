"""Rewrite the AI Leaderboard table in the training deck from the live CSV.

Cell text only: fonts, fills, column widths and the table style are left
exactly as the deck author set them. The table is found by its header text
rather than by slide number, so inserting slides earlier in the deck cannot
silently retarget the update.

    python deck/update_slide.py --deck AI_Training.pptx --dry-run
    python deck/update_slide.py --deck AI_Training.pptx

Writes a copy by default. --in-place overwrites, and only after the same
checks pass.
"""

from __future__ import annotations

import argparse
import copy
import csv
import datetime as dt
import re
import shutil
import sys
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.table import Table

ROOT = Path(__file__).resolve().parent.parent
DEFAULT_CONFIG = Path(__file__).parent / "config.yaml"
DEFAULT_CSV = ROOT / "data" / "leaderboard_latest.csv"

# Header cells carry footnote markers ("GPQA¹", "AIME 2025²"); strip anything
# that is not a letter, digit, space or the $/M punctuation before matching.
_NORMALISE = re.compile(r"[^A-Za-z0-9 $/.]+")


def normalise_header(text: str) -> str:
    return _NORMALISE.sub("", text).strip().lower()


# --------------------------------------------------------------------------
# data
# --------------------------------------------------------------------------

def load_rows(csv_path: Path) -> list[dict[str, Any]]:
    with csv_path.open(encoding="utf-8") as fh:
        return list(csv.DictReader(fh))


def to_number(raw: str) -> float | None:
    if raw is None or raw == "":
        return None
    try:
        return float(raw)
    except ValueError:
        return None


def select_rows(rows: list[dict], selection: dict) -> tuple[list[dict], list[str]]:
    """Pick the rows the slide prints. Returns (rows, warnings)."""
    warnings: list[str] = []
    if selection.get("mode") == "top_n":
        ranked = [r for r in rows if to_number(r.get("gpqa")) is not None]
        ranked.sort(key=lambda r: -to_number(r["gpqa"]))
        return ranked[: int(selection.get("top_n", 8))], warnings

    by_name = {r["model"]: r for r in rows}
    chosen: list[dict] = []
    for name in selection.get("shortlist", []):
        row = by_name.get(name)
        if row is None:
            # A retired model must not shift every later row up by one - that
            # would put the wrong numbers beside the right names.
            warnings.append(f"{name!r} is no longer on the leaderboard; row left blank")
            chosen.append({"model": name, "_missing": True})
        else:
            chosen.append(row)
    return chosen, warnings


# --------------------------------------------------------------------------
# locating the table
# --------------------------------------------------------------------------

def iter_tables(prs: Presentation):
    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_table:
                yield index, slide, shape.table


def find_table(prs: Presentation, cfg: dict) -> tuple[int, Any, Table]:
    wanted = [normalise_header(h) for h in cfg["header_contains"]]
    hint = cfg.get("slide_hint")
    candidates = []
    for index, slide, table in iter_tables(prs):
        headers = [normalise_header(c.text) for c in table.rows[0].cells]
        if all(any(w in h for h in headers) for w in wanted):
            candidates.append((index, slide, table))
    if not candidates:
        raise SystemExit(
            f"No table found whose header row contains {cfg['header_contains']}.\n"
            "Check deck/config.yaml -> table.header_contains against the real deck."
        )
    if hint:
        for cand in candidates:
            if cand[0] == hint:
                return cand
        print(
            f"note: no matching table on slide {hint}; "
            f"using slide {candidates[0][0]} instead"
        )
    if len(candidates) > 1:
        print(f"note: {len(candidates)} matching tables; using slide {candidates[0][0]}")
    return candidates[0]


def map_columns(table: Table, mapping: dict[str, str]) -> dict[int, str]:
    """Header column index -> CSV field, by matching header text."""
    headers = [normalise_header(c.text) for c in table.rows[0].cells]
    resolved: dict[int, str] = {}
    for label, field in mapping.items():
        key = normalise_header(label)
        for i, header in enumerate(headers):
            if header.startswith(key) or key in header:
                resolved[i] = field
                break
        else:
            print(f"note: header {label!r} not found in the table; column skipped")
    return resolved


# --------------------------------------------------------------------------
# writing
# --------------------------------------------------------------------------

def set_cell_text(cell, text: str) -> None:
    """Replace a cell's text while keeping its formatting.

    The first run holds the cell's font. Writing into it and dropping the
    remaining runs preserves size, weight and colour; assigning to
    ``cell.text`` would reset all of it to the theme default.
    """
    para = cell.text_frame.paragraphs[0]
    if not para.runs:
        para.add_run()
    para.runs[0].text = text
    for run in para.runs[1:]:
        run._r.getparent().remove(run._r)
    for extra in cell.text_frame.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)


def render(field: str, row: dict, formats: dict, missing: str,
           display_names: dict[str, str] | None = None) -> str:
    names = display_names or {}
    if field == "model":
        name = row.get("model", "")
        return names.get(name, name)
    if row.get("_missing"):
        return missing
    raw = row.get(field, "")
    fmt = formats.get(field)
    if fmt is None:
        return str(raw)
    value = to_number(raw)
    return missing if value is None else fmt.format(value)


def clone_row(table: Table, template_index: int) -> None:
    """Append a copy of an existing row, inheriting its formatting."""
    tbl = table._tbl
    new_row = copy.deepcopy(tbl.tr_lst[template_index])
    tbl.append(new_row)


def fit_row_count(table: Table, needed: int) -> list[str]:
    """Grow or shrink the body so it holds exactly `needed` rows."""
    notes: list[str] = []
    body = len(table.rows) - 1
    while len(table.rows) - 1 < needed:
        clone_row(table, len(table.rows) - 1)
    removed = 0
    while len(table.rows) - 1 > needed:
        tbl = table._tbl
        tbl.remove(tbl.tr_lst[-1])
        removed += 1
    if len(table.rows) - 1 != body:
        notes.append(f"table body resized from {body} to {needed} rows")
    return notes


def update_footnote(slide, cfg: dict, as_of: dt.date) -> str | None:
    pattern = re.compile(cfg["match"])
    replacement = cfg["replace"].format(date=as_of)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if pattern.search(run.text):
                    before = run.text
                    run.text = pattern.sub(replacement, run.text)
                    return f"{before!r} -> {run.text!r}"
    return None


# --------------------------------------------------------------------------

def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--config", type=Path, default=DEFAULT_CONFIG)
    ap.add_argument("--csv", type=Path, default=DEFAULT_CSV)
    ap.add_argument("--deck", type=Path, help="overrides deck: in the config")
    ap.add_argument("--out", type=Path, help="overrides output: in the config")
    ap.add_argument("--in-place", action="store_true", help="overwrite the deck")
    ap.add_argument("--dry-run", action="store_true", help="report, write nothing")
    args = ap.parse_args()

    cfg = yaml.safe_load(args.config.read_text(encoding="utf-8"))
    deck = args.deck or (ROOT / cfg["deck"])
    if not deck.exists():
        raise SystemExit(f"Deck not found: {deck}")
    if not args.csv.exists():
        raise SystemExit(
            f"No data at {args.csv}. Run the refresh workflow, or pull the branch "
            "the bot commits to."
        )

    rows = load_rows(args.csv)
    chosen, warnings = select_rows(rows, cfg["selection"])
    as_of_raw = next((r.get("as_of_date") for r in rows if r.get("as_of_date")), None)
    as_of = dt.date.fromisoformat(as_of_raw) if as_of_raw else dt.date.today()

    prs = Presentation(str(deck))
    slide_no, slide, table = find_table(prs, cfg["table"])
    columns = map_columns(table, cfg["columns"])
    if not columns:
        raise SystemExit("No table columns could be mapped; check deck/config.yaml.")

    print(f"deck   : {deck}")
    print(f"table  : slide {slide_no}, {len(table.rows) - 1} body rows, "
          f"{len(columns)} mapped columns")
    print(f"data   : {args.csv} ({len(rows)} models, as of {as_of})")
    for warning in warnings:
        print(f"WARNING: {warning}")

    for note in fit_row_count(table, len(chosen)):
        print(f"note   : {note}")

    formats, missing = cfg["formats"], cfg["missing"]
    display_names = cfg.get("display_names") or {}
    for offset, row in enumerate(chosen, start=1):
        cells = table.rows[offset].cells
        for col, field in columns.items():
            if col >= len(cells):
                continue
            text = render(field, row, formats, missing, display_names)
            if args.dry_run:
                current = cells[col].text.strip()
                if current != text:
                    print(f"  r{offset}c{col}: {current!r} -> {text!r}")
            else:
                set_cell_text(cells[col], text)

    change = update_footnote(slide, cfg["footnote"], as_of)
    print(f"footnote: {change}" if change else
          "footnote: no 'As of ...' text found (check footnote.match)")

    if args.dry_run:
        print("\ndry run - nothing written")
        return 0

    target = deck if args.in_place else (args.out or ROOT / cfg["output"])
    if args.in_place:
        backup = deck.with_suffix(deck.suffix + ".bak")
        shutil.copy2(deck, backup)
        print(f"backup : {backup}")
    prs.save(str(target))
    print(f"written: {target}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
