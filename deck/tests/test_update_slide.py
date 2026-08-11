"""End-to-end tests for the deck updater against a synthetic slide 144.

The real deck is confidential and not in this repo, so the fixture rebuilds
its essentials: a table whose headers carry footnote markers, a body row with
deliberate formatting to prove formatting survives, and an "As of ..."
footnote.
"""

from __future__ import annotations

import csv
import subprocess
import sys
from pathlib import Path

import pytest
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt

SCRIPT = Path(__file__).resolve().parents[1] / "update_slide.py"
CONFIG = Path(__file__).resolve().parents[1] / "config.yaml"

HEADERS = ["Organization", "Model", "GPQA¹", "AIME 2025²", "Input $/M", "Output $/M"]

CSV_ROWS = [
    # organization, model, gpqa, aime, in, out
    ("Google",    "Gemini 3 Pro",    "91.9", "100.0", "",     ""),
    ("xAI",       "Grok-4 Heavy",    "88.4", "100.0", "",     ""),
    ("OpenAI",    "GPT-5.1",         "88.1", "94.0",  "1.25", "10.0"),
    ("OpenAI",    "GPT-5.1 Instant", "88.1", "94.0",  "1.25", "10.0"),
    ("OpenAI",    "GPT-5.1 Think.",  "88.1", "94.0",  "1.25", "10.0"),
    ("xAI",       "Grok-4",          "87.5", "91.7",  "3.0",  "15.0"),
    ("Anthropic", "Claude Opus 4.5", "87.0", "",      "5.0",  "25.0"),
    ("Google",    "Gemini 2.5 Pro",  "83.0", "83.0",  "1.25", "10.0"),
]


@pytest.fixture
def workspace(tmp_path: Path) -> dict:
    csv_path = tmp_path / "leaderboard_latest.csv"
    with csv_path.open("w", newline="", encoding="utf-8") as fh:
        w = csv.writer(fh, lineterminator="\n")
        w.writerow([
            "organization", "model", "gpqa", "aime_2025", "input_usd_per_m",
            "output_usd_per_m", "context_window", "is_open_source", "license",
            "as_of_date", "source_url",
        ])
        for org, model, gpqa, aime, pin, pout in CSV_ROWS:
            w.writerow([org, model, gpqa, aime, pin, pout, "", "False",
                        "proprietary", "2026-08-11", "https://llm-stats.com"])

    cfg = yaml.safe_load(CONFIG.read_text(encoding="utf-8"))
    cfg["table"]["slide_hint"] = 1
    cfg_path = tmp_path / "config.yaml"
    cfg_path.write_text(yaml.safe_dump(cfg), encoding="utf-8")
    return {"csv": csv_path, "config": cfg_path, "tmp": tmp_path}


def build_deck(path: Path, body_rows: int = 8) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(
        body_rows + 1, len(HEADERS), Inches(0.4), Inches(1.0), Inches(9), Inches(4)
    )
    table = shape.table
    for col, text in enumerate(HEADERS):
        table.rows[0].cells[col].text = text
    for row in range(1, body_rows + 1):
        for col in range(len(HEADERS)):
            cell = table.rows[row].cells[col]
            cell.text = "placeholder"
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(11)
            run.font.bold = (col == 1)   # model column is bold in the deck

    box = slide.shapes.add_textbox(Inches(6.5), Inches(5.4), Inches(3), Inches(0.4))
    box.text_frame.paragraphs[0].add_run().text = "As of december 2025"
    prs.save(str(path))


def run(workspace: dict, deck: Path, *extra: str) -> subprocess.CompletedProcess:
    return subprocess.run(
        [sys.executable, str(SCRIPT),
         "--config", str(workspace["config"]),
         "--csv", str(workspace["csv"]),
         "--deck", str(deck), *extra],
        capture_output=True, text=True, check=False,
    )


def read_table(path: Path):
    table = next(s.table for s in Presentation(str(path)).slides[0].shapes if s.has_table)
    return [[c.text for c in row.cells] for row in table.rows]


def test_writes_every_shortlist_row(workspace):
    deck = workspace["tmp"] / "deck.pptx"
    build_deck(deck)
    out = workspace["tmp"] / "out.pptx"
    res = run(workspace, deck, "--out", str(out))
    assert res.returncode == 0, res.stdout + res.stderr

    grid = read_table(out)
    assert grid[0][2].startswith("GPQA")            # header untouched
    assert grid[1] == ["Google", "Gemini 3 Pro", "91.9 %", "100.0 %", "n/a", "n/a"]
    assert grid[3] == ["OpenAI", "GPT-5.1", "88.1 %", "94.0 %", "1.25", "10.00"]
    assert grid[7][3] == "n/a"                      # Opus 4.5 has no AIME score
    assert grid[8][1] == "Gemini 2.5 Pro"


def test_withdrawn_price_never_becomes_zero(workspace):
    deck = workspace["tmp"] / "deck.pptx"
    build_deck(deck)
    out = workspace["tmp"] / "out.pptx"
    run(workspace, deck, "--out", str(out))
    flat = [cell for row in read_table(out)[1:] for cell in row]
    assert "n/a" in flat
    assert "0.00" not in flat


def test_cell_formatting_survives(workspace):
    deck = workspace["tmp"] / "deck.pptx"
    build_deck(deck)
    out = workspace["tmp"] / "out.pptx"
    run(workspace, deck, "--out", str(out))

    table = next(s.table for s in Presentation(str(out)).slides[0].shapes if s.has_table)
    run0 = table.rows[1].cells[1].text_frame.paragraphs[0].runs[0]
    assert run0.font.bold is True
    assert run0.font.size == Pt(11)


def test_footnote_date_is_rewritten(workspace):
    deck = workspace["tmp"] / "deck.pptx"
    build_deck(deck)
    out = workspace["tmp"] / "out.pptx"
    run(workspace, deck, "--out", str(out))
    texts = [
        s.text_frame.text
        for s in Presentation(str(out)).slides[0].shapes
        if s.has_text_frame
    ]
    assert any("August 2026" in t for t in texts), texts
    assert not any("december 2025" in t for t in texts)


def test_table_grows_to_fit_the_shortlist(workspace):
    deck = workspace["tmp"] / "deck.pptx"
    build_deck(deck, body_rows=5)
    out = workspace["tmp"] / "out.pptx"
    res = run(workspace, deck, "--out", str(out))
    assert res.returncode == 0, res.stdout + res.stderr
    grid = read_table(out)
    assert len(grid) == 9                      # header + 8
    assert grid[8][1] == "Gemini 2.5 Pro"


def test_retired_model_keeps_its_row_instead_of_shifting(workspace):
    """A missing model must not pull later rows up under the wrong names."""
    cfg = yaml.safe_load(workspace["config"].read_text())
    cfg["selection"]["shortlist"] = ["Gemini 3 Pro", "Model That Retired", "GPT-5.1"]
    workspace["config"].write_text(yaml.safe_dump(cfg))

    deck = workspace["tmp"] / "deck.pptx"
    build_deck(deck)
    out = workspace["tmp"] / "out.pptx"
    res = run(workspace, deck, "--out", str(out))
    assert "no longer on the leaderboard" in res.stdout

    grid = read_table(out)
    assert grid[1][1] == "Gemini 3 Pro"
    assert grid[2][1] == "Model That Retired"
    assert grid[2][2] == "n/a"
    assert grid[3][1] == "GPT-5.1"      # still beside its own numbers


def test_dry_run_writes_nothing(workspace):
    deck = workspace["tmp"] / "deck.pptx"
    build_deck(deck)
    before = deck.read_bytes()
    res = run(workspace, deck, "--dry-run")
    assert res.returncode == 0
    assert "dry run" in res.stdout
    assert deck.read_bytes() == before
    assert not (workspace["tmp"] / "out.pptx").exists()


def test_missing_table_fails_loudly(workspace):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    deck = workspace["tmp"] / "empty.pptx"
    prs.save(str(deck))
    res = run(workspace, deck, "--out", str(workspace["tmp"] / "out.pptx"))
    assert res.returncode != 0
    assert "No table found" in (res.stdout + res.stderr)


def test_display_name_shortens_the_label_without_changing_the_match(workspace):
    """The slide prints "GPT-5.1 Think."; the data says "GPT-5.1 Thinking"."""
    import csv as _csv
    with workspace["csv"].open("a", newline="", encoding="utf-8") as fh:
        _csv.writer(fh, lineterminator="\n").writerow([
            "OpenAI", "GPT-5.1 Thinking", "88.1", "94.0", "", "",
            "", "False", "proprietary", "2026-08-11", "https://llm-stats.com",
        ])
    cfg = yaml.safe_load(workspace["config"].read_text())
    cfg["selection"]["shortlist"] = ["GPT-5.1 Thinking"]
    cfg["display_names"] = {"GPT-5.1 Thinking": "GPT-5.1 Think."}
    workspace["config"].write_text(yaml.safe_dump(cfg))

    deck = workspace["tmp"] / "deck.pptx"
    build_deck(deck)
    out = workspace["tmp"] / "out.pptx"
    res = run(workspace, deck, "--out", str(out))
    assert res.returncode == 0, res.stdout + res.stderr
    assert "no longer on the leaderboard" not in res.stdout   # matched on the real name

    row = read_table(out)[1]
    assert row[1] == "GPT-5.1 Think."     # short label on the slide
    assert row[2] == "88.1 %"             # real model's numbers
    assert row[5] == "n/a"                # price withdrawn upstream
